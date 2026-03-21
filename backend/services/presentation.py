"""Сервис генерации .pptx из RAG-контекста."""

from __future__ import annotations

import json
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from openai import OpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from backend.core.config import settings
from backend.utils.embeddings import normalize_base_url, resolve_api_key


def _sanitize_filename(name: str) -> str:
    safe = re.sub(r"[^a-zA-Z0-9._-]+", "_", name).strip("._")
    return safe or "presentation"


def _extract_json(raw: str) -> Dict[str, Any]:
    text = raw.strip()
    if text.startswith("```"):
        text = text.strip("`")
        if text.startswith("json"):
            text = text[4:].strip()
    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise ValueError("LLM response does not contain JSON object")
    return json.loads(text[start : end + 1])


def _build_context(results: List[Dict[str, Any]]) -> str:
    lines: List[str] = []
    for idx, item in enumerate(results, 1):
        lines.append(
            "\n".join(
                [
                    f"[source#{idx}]",
                    f"source_name={item.get('source_name', 'unknown')}",
                    f"chunk_index={item.get('chunk_index', -1)}",
                    f"text={str(item.get('text', '')).strip()}",
                ]
            )
        )
    return "\n\n".join(lines)


def _request_slide_plan(
    query: str,
    results: List[Dict[str, Any]],
    max_slides: int,
    model: Optional[str],
) -> Dict[str, Any]:
    llm = settings.llm
    api_key = resolve_api_key(llm.api_key, ["HACKAI_API_KEY", "OPENAI_API_KEY"])
    if not api_key:
        raise RuntimeError("LLM API key required")

    client = OpenAI(
        api_key=api_key,
        base_url=normalize_base_url(llm.base_url),
        timeout=180.0,
        max_retries=2,
    )
    context = _build_context(results)
    max_bullets = settings.presentation.max_bullets_per_slide
    prompt = (
        "Собери структуру презентации на русском языке только по данным из контекста.\n"
        "Выведи строго JSON без markdown.\n"
        "Формат:\n"
        "{\n"
        '  "title": "строка",\n'
        '  "slides": [\n'
        "    {\n"
        '      "title": "строка",\n'
        '      "bullets": ["пункт 1", "пункт 2"],\n'
        '      "source_ids": [1, 2]\n'
        "    }\n"
        "  ]\n"
        "}\n"
        f"Ограничения: не более {max_slides} слайдов и не более {max_bullets} bullets на слайд. "
        "Каждый пункт 8-20 слов. Если данных мало, сделай меньше слайдов."
    )
    response = client.responses.create(
        model=model or llm.model,
        input=[
            {
                "role": "system",
                "content": (
                    "Ты аналитик RAG-системы. Формируй только проверяемые тезисы "
                    "из исходного контекста. Не добавляй сведения вне контекста."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Запрос пользователя:\n{query}\n\n"
                    f"Контекст:\n{context}\n\n{prompt}"
                ),
            },
        ],
        temperature=0.2,
    )
    return _extract_json(response.output_text)


def _to_int_ids(raw: Any, max_id: int) -> List[int]:
    if not isinstance(raw, list):
        return []
    result: List[int] = []
    for val in raw:
        try:
            idx = int(val)
        except (TypeError, ValueError):
            continue
        if 1 <= idx <= max_id and idx not in result:
            result.append(idx)
    return result


def _normalize_plan(
    plan: Dict[str, Any],
    results: List[Dict[str, Any]],
    max_slides: int,
) -> Dict[str, Any]:
    max_bullets = settings.presentation.max_bullets_per_slide
    title = str(plan.get("title") or "Презентация по материалам")
    raw_slides = plan.get("slides")

    slides: List[Dict[str, Any]] = []
    if isinstance(raw_slides, list):
        for row in raw_slides[:max_slides]:
            if not isinstance(row, dict):
                continue
            slide_title = str(row.get("title") or "").strip()
            bullets_raw = row.get("bullets")
            if not slide_title or not isinstance(bullets_raw, list):
                continue
            bullets = [str(item).strip() for item in bullets_raw if str(item).strip()]
            bullets = bullets[:max_bullets]
            if not bullets:
                continue
            slides.append(
                {
                    "title": slide_title,
                    "bullets": bullets,
                    "source_ids": _to_int_ids(row.get("source_ids"), len(results)),
                }
            )

    if not slides:
        fallback = []
        for item in results[:max_slides]:
            text = str(item.get("text", "")).strip()
            if not text:
                continue
            fallback.append(
                {
                    "title": str(item.get("source_name", "Источник")),
                    "bullets": [text[:260]],
                    "source_ids": [],
                }
            )
        slides = fallback or [
            {"title": "Недостаточно данных", "bullets": ["Релевантные источники не найдены"], "source_ids": []}
        ]

    return {"title": title, "slides": slides}


def _render_pptx(plan: Dict[str, Any], results: List[Dict[str, Any]], target_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Pt(960)   # 13.333"
    prs.slide_height = Pt(540)  # 7.5" (16:9)

    # Титульный слайд
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_fill = title_slide.background.fill
    title_fill.solid()
    title_fill.fore_color.rgb = RGBColor(16, 185, 129)
    title_box = title_slide.shapes.add_textbox(Pt(72), Pt(180), Pt(816), Pt(180))
    title_tf = title_box.text_frame
    title_tf.clear()
    title_tf.word_wrap = True
    title_title = title_tf.paragraphs[0]
    title_title.text = str(plan["title"])
    title_title.font.name = "Inter"
    title_title.font.bold = True
    title_title.font.color.rgb = RGBColor(255, 255, 255)
    title_title.font.size = Pt(56)
    title_title.alignment = PP_ALIGN.CENTER

    # Контентные слайды
    for item in plan["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        content_fill = slide.background.fill
        content_fill.solid()
        content_fill.fore_color.rgb = RGBColor(15, 23, 42)

        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Pt(0), Pt(0), Pt(960), Pt(20)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(99, 102, 241)
        accent.line.fill.background()

        slide.shapes.title.text = str(item["title"])
        slide_title_p = slide.shapes.title.text_frame.paragraphs[0]
        slide_title_p.font.name = "Inter"
        slide_title_p.font.bold = True
        slide_title_p.font.size = Pt(34)
        slide_title_p.font.color.rgb = RGBColor(255, 255, 255)

        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for idx, bullet in enumerate(item["bullets"]):
            paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.name = "Inter"
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = RGBColor(203, 213, 225)
            paragraph.space_after = Pt(10)

    target_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(target_path))


def generate_presentation(  # pylint: disable=too-many-arguments,too-many-positional-arguments
    query: str,
    results: List[Dict[str, Any]],
    *,
    model: Optional[str] = None,
    max_slides: Optional[int] = None,
    filename_prefix: Optional[str] = None,
) -> Tuple[Path, Dict[str, Any]]:
    """Генерирует pptx на основе результатов retrieve и возвращает путь + метаданные."""
    if not results:
        raise ValueError("No retrieval results")

    cfg = settings.presentation
    requested = max_slides if max_slides is not None else cfg.default_max_slides
    bounded_max_slides = max(2, min(requested, cfg.max_slides_limit))

    plan_raw = _request_slide_plan(query, results, bounded_max_slides, model)
    plan = _normalize_plan(plan_raw, results, bounded_max_slides)

    stamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    base = _sanitize_filename(filename_prefix or query[:80] or "presentation")
    path = Path(cfg.output_dir).resolve() / f"{base}_{stamp}.pptx"
    _render_pptx(plan, results, path)

    metadata = {
        "title": plan["title"],
        "slides": len(plan["slides"]) + 1,  # + титульный слайд
        "path": str(path),
    }
    return path, metadata
