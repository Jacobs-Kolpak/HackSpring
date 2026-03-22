from __future__ import annotations

import json
import re
from copy import deepcopy
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from backend.core.config import settings
from backend.utils.llm import generate_text

DEFAULT_PRESENTATION_STYLE: Dict[str, Any] = {
    "title_background": "#10B981",
    "title_text": "#FFFFFF",
    "slide_background": "#0F172A",
    "accent": "#6366F1",
    "slide_title_text": "#FFFFFF",
    "bullet_text": "#CBD5E1",
    "font_family": "Inter",
    "title_font_size": 56,
    "slide_title_font_size": 34,
    "bullet_font_size": 24,
}
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _sanitize_filename(name: str) -> str:
    safe = re.sub(r"[^a-zA-Z0-9._-]+", "_", name).strip("._")
    return safe or "presentation"


def _extract_json(raw: str) -> Dict[str, Any]:
    text = raw.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end <= start:
        start = text.find("[")
        end = text.rfind("]")
        if start == -1 or end == -1 or end <= start:
            raise ValueError("LLM response does not contain JSON")
    fragment = text[start : end + 1]
    try:
        return json.loads(fragment)
    except json.JSONDecodeError:
        cleaned = re.sub(r",\s*([}\]])", r"\1", fragment)
        cleaned = re.sub(r"'", '"', cleaned)
        cleaned = re.sub(r"\n", " ", cleaned)
        return json.loads(cleaned)


def _safe_extract_json(raw: str, model: Optional[str]) -> Dict[str, Any]:
    try:
        return _extract_json(raw)
    except Exception:
        repair_prompt = (
            "Исправь невалидный JSON. Верни строго валидный JSON-объект "
            "без markdown и комментариев. Не меняй смысл, только синтаксис.\n\n"
            f"{raw}"
        )
        repaired = generate_text(
            repair_prompt,
            system=(
                "Ты валидатор JSON. Возвращай только валидный JSON-объект, "
                "ничего кроме JSON."
            ),
            model=model,
            temperature=0.0,
            max_tokens=4096,
        )
        return _extract_json(repaired)


def _build_context(results: List[Dict[str, Any]]) -> str:
    lines: List[str] = []
    for idx, item in enumerate(results, 1):
        lines.append(
            "\n".join([
                f"[source#{idx}]",
                f"source_name={item.get('source_name', 'unknown')}",
                f"chunk_index={item.get('chunk_index', -1)}",
                f"text={str(item.get('text', '')).strip()}",
            ])
        )
    return "\n\n".join(lines)


def _request_slide_plan(
    query: str,
    results: List[Dict[str, Any]],
    max_slides: int,
    model: Optional[str],
) -> Dict[str, Any]:
    context = _build_context(results)
    max_bullets = settings.presentation.max_bullets_per_slide
    prompt = (
        "Собери структуру презентации на русском языке только по данным из контекста.\n"
        "Выведи строго JSON без markdown.\n"
        "Формат:\n"
        "{\n"
        '  "title": "строка",\n'
        '  "slides": [\n'
        "    {\n"
        '      "title": "строка",\n'
        '      "bullets": ["пункт 1", "пункт 2"],\n'
        '      "source_ids": [1, 2]\n'
        "    }\n"
        "  ]\n"
        "}\n"
        f"Ограничения: не более {max_slides} слайдов и не более {max_bullets} bullets на слайд. "
        "Каждый пункт 8-20 слов. Если данных мало, сделай меньше слайдов.\n\n"
        f"Запрос пользователя:\n{query}\n\n"
        f"Контекст:\n{context}"
    )
    raw = generate_text(
        prompt,
        system=(
            "Ты аналитик RAG-системы. Формируй только проверяемые тезисы "
            "из исходного контекста. Не добавляй сведения вне контекста. "
            "Верни только JSON."
        ),
        model=model,
        temperature=0.2,
        max_tokens=4096,
    )
    try:
        return _safe_extract_json(raw, model)
    except Exception:
        # Позволяем пайплайну перейти на fallback в _normalize_plan
        return {"title": "Презентация по материалам", "slides": []}


def _to_int_ids(raw: Any, max_id: int) -> List[int]:
    if not isinstance(raw, list):
        return []
    result: List[int] = []
    for val in raw:
        try:
            idx = int(val)
        except (TypeError, ValueError):
            continue
        if 1 <= idx <= max_id and idx not in result:
            result.append(idx)
    return result


def _normalize_plan(
    plan: Dict[str, Any],
    results: List[Dict[str, Any]],
    max_slides: int,
) -> Dict[str, Any]:
    max_bullets = settings.presentation.max_bullets_per_slide
    title = str(plan.get("title") or "Презентация по материалам")
    raw_slides = plan.get("slides")

    slides: List[Dict[str, Any]] = []
    if isinstance(raw_slides, list):
        for row in raw_slides[:max_slides]:
            if not isinstance(row, dict):
                continue
            slide_title = str(row.get("title") or "").strip()
            bullets_raw = row.get("bullets")
            if not slide_title or not isinstance(bullets_raw, list):
                continue
            bullets = [
                str(item).strip()
                for item in bullets_raw
                if str(item).strip()
            ]
            bullets = bullets[:max_bullets]
            if not bullets:
                continue
            slides.append({
                "title": slide_title,
                "bullets": bullets,
                "source_ids": _to_int_ids(row.get("source_ids"), len(results)),
            })

    if not slides:
        fallback = []
        for item in results[:max_slides]:
            text = str(item.get("text", "")).strip()
            if not text:
                continue
            fallback.append({
                "title": str(item.get("source_name", "Источник")),
                "bullets": [text[:260]],
                "source_ids": [],
            })
        slides = fallback or [{
            "title": "Недостаточно данных",
            "bullets": ["Релевантные источники не найдены"],
            "source_ids": [],
        }]

    return {"title": title, "slides": slides}


def _rgb_to_hex(raw_rgb: Any) -> Optional[str]:
    if raw_rgb is None:
        return None
    value = str(raw_rgb).strip()
    if len(value) != 6 or re.search(r"[^0-9a-fA-F]", value):
        return None
    return f"#{value.upper()}"


def _extract_fill_hex(fill: Any) -> Optional[str]:
    if fill is None:
        return None
    try:
        return _rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        return None


def _extract_text_frame_style(text_frame: Any) -> Dict[str, Any]:
    out: Dict[str, Any] = {}
    if text_frame is None:
        return out

    for paragraph in text_frame.paragraphs:
        fonts = [run.font for run in paragraph.runs] or [paragraph.font]
        for font in fonts:
            if "font_family" not in out:
                name = getattr(font, "name", None)
                if name:
                    out["font_family"] = str(name).strip()
            if "font_size" not in out:
                size = getattr(font, "size", None)
                pt = getattr(size, "pt", None)
                if pt:
                    out["font_size"] = int(round(pt))
            if "font_color" not in out:
                try:
                    color = _rgb_to_hex(font.color.rgb)
                except Exception:
                    color = None
                if color:
                    out["font_color"] = color
            if (
                "font_family" in out
                and "font_size" in out
                and "font_color" in out
            ):
                return out
    return out


def _pick_layout(prs: Presentation, preferred_index: int):
    if len(prs.slide_layouts) == 0:
        raise ValueError("Template has no slide layouts")
    safe_idx = max(0, min(preferred_index, len(prs.slide_layouts) - 1))
    return prs.slide_layouts[safe_idx]


def _remove_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)  # type: ignore[attr-defined] # pylint: disable=protected-access
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)  # type: ignore[attr-defined] # pylint: disable=protected-access


def _remove_first_n_slides(prs: Presentation, count: int) -> None:
    for _ in range(max(0, count)):
        slide_ids = prs.slides._sldIdLst  # type: ignore[attr-defined] # pylint: disable=protected-access
        if len(slide_ids) == 0:
            break
        slide_id = slide_ids[0]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_ids.remove(slide_id)


def _copy_slide_relationships(source_slide: Any, target_slide: Any) -> Dict[str, str]:
    rel_map: Dict[str, str] = {}
    for rel in source_slide.part.rels.values():
        reltype = str(rel.reltype or "")
        if reltype.endswith("/slideLayout") or reltype.endswith("/notesSlide"):
            continue
        if rel.is_external:
            new_rid = target_slide.part.rels.get_or_add_ext_rel(reltype, rel.target_ref)
        else:
            new_rid = target_slide.part.rels.get_or_add(reltype, rel.target_part)
        rel_map[rel.rId] = new_rid
    return rel_map


def _remap_rel_ids(element: Any, rel_map: Dict[str, str]) -> None:
    for node in element.iter():
        for key, value in list(node.attrib.items()):
            if not key.startswith(f"{{{REL_NS}}}"):
                continue
            if value in rel_map:
                node.set(key, rel_map[value])


def _copy_template_visuals(source_slide: Any, target_slide: Any) -> None:
    rel_map = _copy_slide_relationships(source_slide, target_slide)

    source_bg = source_slide._element.cSld.bg
    if source_bg is not None:
        cloned_bg = deepcopy(source_bg)
        _remap_rel_ids(cloned_bg, rel_map)
        if target_slide._element.cSld.bg is not None:
            target_slide._element.cSld._remove_bg()
        target_slide._element.cSld._insert_bg(cloned_bg)

    for shape in source_slide.shapes:
        if getattr(shape, "is_placeholder", False):
            continue
        cloned_shape = deepcopy(shape.element)
        _remap_rel_ids(cloned_shape, rel_map)
        target_slide.shapes._spTree.insert_element_before(cloned_shape, "p:extLst")


def _pick_template_source_slide(index: int, source_slides: List[Any]) -> Any:
    if not source_slides:
        return None
    if index < len(source_slides):
        return source_slides[index]
    if len(source_slides) > 1:
        return source_slides[1]
    return source_slides[0]


def _get_title_text_frame(slide: Any) -> Any:
    title_shape = slide.shapes.title
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        return title_shape.text_frame
    fallback = slide.shapes.add_textbox(Pt(72), Pt(60), Pt(816), Pt(90))
    return fallback.text_frame


def _get_body_text_frame(slide: Any, *, create_if_missing: bool) -> Any:
    title_shape = slide.shapes.title
    allowed_types = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
    maybe_text = getattr(PP_PLACEHOLDER, "TEXT", None)
    if maybe_text is not None:
        allowed_types.add(maybe_text)

    for placeholder in slide.placeholders:
        if not getattr(placeholder, "has_text_frame", False):
            continue
        if title_shape is not None and placeholder == title_shape:
            continue
        try:
            ph_type = placeholder.placeholder_format.type
        except Exception:
            ph_type = None
        if ph_type in allowed_types:
            return placeholder.text_frame

    for shape in slide.shapes:
        if title_shape is not None and shape == title_shape:
            continue
        if getattr(shape, "has_text_frame", False):
            return shape.text_frame

    if not create_if_missing:
        return None

    fallback = slide.shapes.add_textbox(Pt(72), Pt(140), Pt(816), Pt(320))
    return fallback.text_frame


def _extract_paragraph_style(paragraph: Any) -> Dict[str, Any]:
    style: Dict[str, Any] = {
        "level": getattr(paragraph, "level", 0),
        "alignment": getattr(paragraph, "alignment", None),
        "line_spacing": getattr(paragraph, "line_spacing", None),
        "space_before": getattr(paragraph, "space_before", None),
        "space_after": getattr(paragraph, "space_after", None),
    }
    font = paragraph.runs[0].font if paragraph.runs else paragraph.font
    style["font_name"] = getattr(font, "name", None)
    style["font_size"] = getattr(font, "size", None)
    style["font_bold"] = getattr(font, "bold", None)
    style["font_italic"] = getattr(font, "italic", None)
    style["font_underline"] = getattr(font, "underline", None)
    try:
        style["font_color_rgb"] = font.color.rgb
    except Exception:
        style["font_color_rgb"] = None
    return style


def _apply_paragraph_style(paragraph: Any, style: Dict[str, Any]) -> None:
    paragraph.level = style.get("level", 0)
    paragraph.alignment = style.get("alignment")
    paragraph.line_spacing = style.get("line_spacing")
    paragraph.space_before = style.get("space_before")
    paragraph.space_after = style.get("space_after")


def _compact_paragraph_spacing(paragraph: Any, role: str) -> None:
    if role == "body":
        paragraph.line_spacing = 0.92
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(2)
        return

    if role == "slide_title":
        paragraph.line_spacing = 0.96
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        return

    paragraph.line_spacing = 0.95
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)


def _set_paragraph_text_with_style(paragraph: Any, text: str, style: Dict[str, Any]) -> None:
    for run in list(paragraph.runs):
        paragraph._p.remove(run._r)  # type: ignore[attr-defined] # pylint: disable=protected-access
    run = paragraph.add_run()
    run.text = text

    font = run.font
    font.name = style.get("font_name")
    font.size = style.get("font_size")
    font.bold = style.get("font_bold")
    font.italic = style.get("font_italic")
    font.underline = style.get("font_underline")
    color = style.get("font_color_rgb")
    if color is not None:
        try:
            font.color.rgb = color
        except Exception:
            pass


def _compute_text_scale(
    lines: List[str],
    *,
    role: str,
    template_line_count: int,
) -> float:
    non_empty = [line for line in lines if str(line).strip()]
    line_count = max(1, len(non_empty))
    total_chars = sum(len(str(line)) for line in non_empty)
    avg_len = total_chars / line_count if line_count else 0.0

    if role == "main_title":
        char_factor = 1.0 if total_chars <= 42 else (42.0 / total_chars) ** 0.55
        line_factor = 1.0 if line_count <= 2 else (2.0 / line_count) ** 0.32
        return max(0.50, min(1.0, char_factor * line_factor))

    if role == "slide_title":
        char_factor = 1.0 if total_chars <= 32 else (32.0 / total_chars) ** 0.58
        line_factor = 1.0 if line_count <= 2 else (2.0 / line_count) ** 0.36
        return max(0.46, min(1.0, char_factor * line_factor))

    # Body: агрессивнее уменьшаем при большом количестве пунктов/длине строк.
    expected_lines = max(1, template_line_count)
    count_factor = 1.0 if line_count <= expected_lines else (expected_lines / line_count) ** 0.48
    avg_factor = 1.0 if avg_len <= 56 else (56.0 / avg_len) ** 0.42
    return max(0.42, min(1.0, count_factor * avg_factor))


def _scaled_font_size(source_size: Any, scale: float, role: str) -> Any:
    if source_size is None:
        return None
    base_pt = getattr(source_size, "pt", None)
    if base_pt is None:
        try:
            base_pt = float(source_size)
        except (TypeError, ValueError):
            return source_size

    if role == "main_title":
        min_pt, max_pt = 18.0, 42.0
    elif role == "slide_title":
        min_pt, max_pt = 14.0, 30.0
    else:
        min_pt, max_pt = 11.0, 18.0

    new_pt = max(min_pt, min(max_pt, float(base_pt) * scale))
    return Pt(new_pt)


def _replace_text_keep_template_style(
    text_frame: Any,
    lines: List[str],
    reference_text_frame: Any = None,
    *,
    role: str = "body",
) -> None:
    if not lines:
        lines = [""]

    current = list(text_frame.paragraphs)
    if not current:
        return

    source_frame = reference_text_frame if reference_text_frame is not None else text_frame
    source_paragraphs = list(source_frame.paragraphs)
    if source_paragraphs:
        template_styles = [_extract_paragraph_style(p) for p in source_paragraphs]
    else:
        template_styles = [_extract_paragraph_style(p) for p in current]

    scale = _compute_text_scale(
        lines,
        role=role,
        template_line_count=len(template_styles),
    )
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for idx, text in enumerate(lines):
        if idx < len(current):
            paragraph = current[idx]
        else:
            paragraph = text_frame.add_paragraph()
            current.append(paragraph)
        base_style = template_styles[min(idx, len(template_styles) - 1)]
        styled = {**base_style}
        styled["font_size"] = _scaled_font_size(base_style.get("font_size"), scale, role)
        _apply_paragraph_style(paragraph, base_style)
        _set_paragraph_text_with_style(paragraph, str(text), styled)
        _compact_paragraph_spacing(paragraph, role)

    for idx in range(len(lines), len(current)):
        base_style = template_styles[min(idx, len(template_styles) - 1)]
        _apply_paragraph_style(current[idx], base_style)
        _set_paragraph_text_with_style(current[idx], "", base_style)
        _compact_paragraph_spacing(current[idx], role)


def _clear_other_text_frames(slide: Any, keep_frames: List[Any]) -> None:
    keep_ids = {
        id(getattr(frame, "_element", None))
        for frame in keep_frames
        if frame is not None
    }
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text_frame = shape.text_frame
        if id(getattr(text_frame, "_element", None)) in keep_ids:
            continue
        text_frame.clear()


def _merge_style(base: Dict[str, Any], extra: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    merged = {**base}
    if isinstance(extra, dict):
        for key, value in extra.items():
            if value not in (None, ""):
                merged[key] = value
    return merged


def _extract_style_from_template(template_path: Optional[Path]) -> Dict[str, Any]:
    if template_path is None or not template_path.exists():
        return {}

    try:
        prs = Presentation(str(template_path))
    except Exception:
        return {}

    if len(prs.slides) == 0:
        return {}

    style: Dict[str, Any] = {}
    title_slide = prs.slides[0]
    content_slide = prs.slides[1] if len(prs.slides) > 1 else prs.slides[0]

    title_bg = _extract_fill_hex(title_slide.background.fill)
    if title_bg:
        style["title_background"] = title_bg

    slide_bg = _extract_fill_hex(content_slide.background.fill)
    if slide_bg:
        style["slide_background"] = slide_bg

    title_tf = _get_title_text_frame(title_slide)
    title_st = _extract_text_frame_style(title_tf)
    if "font_color" in title_st:
        style["title_text"] = title_st["font_color"]
    if "font_family" in title_st:
        style["font_family"] = title_st["font_family"]
    if "font_size" in title_st:
        style["title_font_size"] = title_st["font_size"]

    content_title_tf = _get_title_text_frame(content_slide)
    content_title_st = _extract_text_frame_style(content_title_tf)
    if "font_color" in content_title_st:
        style["slide_title_text"] = content_title_st["font_color"]
    if "font_size" in content_title_st:
        style["slide_title_font_size"] = content_title_st["font_size"]
    if "font_family" not in style and "font_family" in content_title_st:
        style["font_family"] = content_title_st["font_family"]

    body_tf = _get_body_text_frame(content_slide, create_if_missing=False)
    body_st = _extract_text_frame_style(body_tf)
    if "font_color" in body_st:
        style["bullet_text"] = body_st["font_color"]
    if "font_size" in body_st:
        style["bullet_font_size"] = body_st["font_size"]
    if "font_family" not in style and "font_family" in body_st:
        style["font_family"] = body_st["font_family"]

    for shape in content_slide.shapes:
        fill_hex = _extract_fill_hex(getattr(shape, "fill", None))
        if fill_hex and fill_hex != style.get("slide_background"):
            style["accent"] = fill_hex
            break

    return style


def _parse_hex_color(value: Any, fallback: str) -> RGBColor:
    raw = str(value or "").strip() or fallback
    normalized = raw[1:] if raw.startswith("#") else raw
    if len(normalized) != 6 or re.search(r"[^0-9a-fA-F]", normalized):
        normalized = fallback[1:] if fallback.startswith("#") else fallback
    try:
        return RGBColor.from_string(normalized.upper())
    except ValueError:
        return RGBColor.from_string("000000")


def _normalize_style(style: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    merged = {**DEFAULT_PRESENTATION_STYLE}
    if isinstance(style, dict):
        for key, value in style.items():
            if key in merged and value not in (None, ""):
                merged[key] = value

    font = str(merged.get("font_family") or "Inter").strip()
    merged["font_family"] = font or "Inter"

    def _size(name: str, fallback: int) -> int:
        raw = merged.get(name, fallback)
        try:
            value = int(raw)
        except (TypeError, ValueError):
            return fallback
        return max(10, min(value, 96))

    merged["title_font_size"] = _size("title_font_size", 56)
    merged["slide_title_font_size"] = _size("slide_title_font_size", 34)
    merged["bullet_font_size"] = _size("bullet_font_size", 24)
    return merged


def _render_pptx(  # pylint: disable=too-many-locals
    plan: Dict[str, Any],
    target_path: Path,
    style: Dict[str, Any],
    *,
    template_path: Optional[Path] = None,
    force_custom_colors: bool = True,
) -> None:
    use_template = template_path is not None and template_path.exists()
    original_template_slides = 0
    template_source_slides: List[Any] = []

    if use_template:
        prs = Presentation(str(template_path))
        original_template_slides = len(prs.slides)
        if original_template_slides > 0:
            template_source_slides = [slide for slide in prs.slides]
        else:
            use_template = False

    if not use_template:
        prs = Presentation()
        prs.slide_width = Pt(960)
        prs.slide_height = Pt(540)

    title_bg = _parse_hex_color(style.get("title_background"), "#10B981")
    title_text = _parse_hex_color(style.get("title_text"), "#FFFFFF")
    slide_bg = _parse_hex_color(style.get("slide_background"), "#0F172A")
    accent = _parse_hex_color(style.get("accent"), "#6366F1")
    slide_title_text = _parse_hex_color(style.get("slide_title_text"), "#FFFFFF")
    bullet_text = _parse_hex_color(style.get("bullet_text"), "#CBD5E1")
    font = str(style.get("font_family") or "Inter")
    title_font_size = int(style.get("title_font_size") or 56)
    slide_title_font_size = int(style.get("slide_title_font_size") or 34)
    bullet_font_size = int(style.get("bullet_font_size") or 24)
    preserve_template_text_style = use_template and not force_custom_colors

    if not use_template:
        title_layout = _pick_layout(prs, 0)
        content_layout = _pick_layout(prs, 1 if len(prs.slide_layouts) > 1 else 0)

    if use_template:
        title_source = _pick_template_source_slide(0, template_source_slides)
        title_slide = prs.slides.add_slide(title_source.slide_layout)
        _copy_template_visuals(title_source, title_slide)
    else:
        title_slide = prs.slides.add_slide(title_layout)
    if force_custom_colors:
        title_fill = title_slide.background.fill
        title_fill.solid()
        title_fill.fore_color.rgb = title_bg

    title_tf = _get_title_text_frame(title_slide)
    if preserve_template_text_style:
        title_ref_tf = _get_title_text_frame(title_source) if use_template else None
        _replace_text_keep_template_style(
            title_tf,
            [str(plan["title"])],
            title_ref_tf,
            role="main_title",
        )
        title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        title_tf.clear()
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = str(plan["title"])
        title_p.font.name = font
        title_p.font.bold = True
        title_p.font.color.rgb = title_text
        title_p.font.size = Pt(title_font_size)
        title_p.alignment = PP_ALIGN.CENTER
        _compact_paragraph_spacing(title_p, "main_title")
    if use_template:
        _clear_other_text_frames(title_slide, [title_tf])

    for source_idx, item in enumerate(plan["slides"], start=1):
        if use_template:
            content_source = _pick_template_source_slide(source_idx, template_source_slides)
            slide = prs.slides.add_slide(content_source.slide_layout)
            _copy_template_visuals(content_source, slide)
        else:
            slide = prs.slides.add_slide(content_layout)
        if force_custom_colors:
            bg_fill = slide.background.fill
            bg_fill.solid()
            bg_fill.fore_color.rgb = slide_bg

        if force_custom_colors and not use_template:
            accent_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, Pt(0), Pt(0), Pt(960), Pt(20)
            )
            accent_shape.fill.solid()
            accent_shape.fill.fore_color.rgb = accent
            accent_shape.line.fill.background()

        slide_title_tf = _get_title_text_frame(slide)
        if preserve_template_text_style:
            title_ref_tf = _get_title_text_frame(content_source) if use_template else None
            _replace_text_keep_template_style(
                slide_title_tf,
                [str(item["title"])],
                title_ref_tf,
                role="slide_title",
            )
        else:
            slide_title_tf.clear()
            slide_title_p = slide_title_tf.paragraphs[0]
            slide_title_p.text = str(item["title"])
            slide_title_p.font.name = font
            slide_title_p.font.bold = True
            slide_title_p.font.size = Pt(slide_title_font_size)
            slide_title_p.font.color.rgb = slide_title_text
            _compact_paragraph_spacing(slide_title_p, "slide_title")

        body = _get_body_text_frame(slide, create_if_missing=True)
        if preserve_template_text_style:
            body_ref_tf = (
                _get_body_text_frame(content_source, create_if_missing=False)
                if use_template
                else None
            )
            _replace_text_keep_template_style(
                body,
                [str(bullet) for bullet in item["bullets"]],
                body_ref_tf,
                role="body",
            )
        else:
            body.clear()
            for idx, bullet in enumerate(item["bullets"]):
                paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 0
                paragraph.font.name = font
                paragraph.font.size = Pt(bullet_font_size)
                paragraph.font.color.rgb = bullet_text
                _compact_paragraph_spacing(paragraph, "body")
        if use_template:
            _clear_other_text_frames(slide, [slide_title_tf, body])

    if use_template and original_template_slides > 0:
        _remove_first_n_slides(prs, original_template_slides)

    target_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(target_path))


def generate_presentation(
    query: str,
    results: List[Dict[str, Any]],
    *,
    model: Optional[str] = None,
    max_slides: Optional[int] = None,
    filename_prefix: Optional[str] = None,
    style: Optional[Dict[str, Any]] = None,
    template_path: Optional[Path] = None,
) -> Tuple[Path, Dict[str, Any]]:
    if not results:
        raise ValueError("No retrieval results")

    cfg = settings.presentation
    requested = max_slides if max_slides is not None else cfg.default_max_slides
    bounded = max(2, min(requested, cfg.max_slides_limit))
    template_style = _extract_style_from_template(template_path)
    style_cfg = _normalize_style(
        _merge_style(
            _merge_style(DEFAULT_PRESENTATION_STYLE, template_style),
            style,
        )
    )

    plan_raw = _request_slide_plan(query, results, bounded, model)
    plan = _normalize_plan(plan_raw, results, bounded)

    stamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    base = _sanitize_filename(filename_prefix or query[:80] or "presentation")
    path = Path(cfg.output_dir).resolve() / f"{base}_{stamp}.pptx"

    template_used = template_path is not None and template_path.exists()
    _render_pptx(
        plan,
        path,
        style_cfg,
        template_path=template_path if template_used else None,
        force_custom_colors=not template_used or bool(style),
    )

    metadata = {
        "title": plan["title"],
        "slides_count": len(plan["slides"]) + 1,
        "path": str(path),
        "style": style_cfg,
        "template_used": template_used,
        "template_name": template_path.name if template_used else None,
    }
    return path, metadata
